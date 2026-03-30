import re
import tempfile
import os
from io import BytesIO
from bs4 import BeautifulSoup
import yt_dlp

import fitz  # PyMuPDF
import httpx
from docx import Document
from pptx import Presentation
from groq import Groq
from app.config import settings


class DocumentParser:
    """Unified parser for ALL formats."""

    def parse_pdf(self, file_bytes: bytes) -> list[dict]:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        for page_num, page in enumerate(doc):
            text = page.get_text("text").strip()
            if text:
                pages.append(
                    {
                        "content": text,
                        "page_number": page_num + 1,
                        "metadata": {"type": "pdf"},
                    }
                )
        doc.close()
        return pages

    def parse_pptx(self, file_bytes: bytes) -> list[dict]:
        prs = Presentation(BytesIO(file_bytes))
        slides = []
        for slide_num, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            texts.append(p.text.strip())
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                if slide.notes_slide.notes_text_frame.text.strip():
                    texts.append(
                        f"[Notes]: {slide.notes_slide.notes_text_frame.text.strip()}"
                    )
            content = "\n".join(texts).strip()
            if content:
                slides.append(
                    {
                        "content": content,
                        "page_number": slide_num + 1,
                        "metadata": {"type": "slide"},
                    }
                )
        return slides

    def parse_docx(self, file_bytes: bytes) -> list[dict]:
        doc = Document(BytesIO(file_bytes))
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return [
            {"content": full_text, "page_number": None, "metadata": {"type": "docx"}}
        ]

    def parse_text(self, text: str) -> list[dict]:
        return [
            {"content": text.strip(), "page_number": None, "metadata": {"type": "text"}}
        ]

    def parse_website(self, url: str) -> list[dict]:
        headers = {"User-Agent": "Mozilla/5.0"}
        with httpx.Client(follow_redirects=True, timeout=15) as client:
            response = client.get(url, headers=headers)
            response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()

        text = soup.get_text(separator="\n", strip=True)
        text = re.sub(r"\n\s*\n", "\n\n", text)
        return [
            {
                "content": text.strip(),
                "page_number": None,
                "metadata": {"type": "website", "url": url},
            }
        ]

    def parse_audio_video(self, file_bytes: bytes, filename: str) -> list[dict]:
        ext = filename.split(".")[-1].lower() if filename else "webm"
        if ext not in ["mp3", "mp4", "mpeg", "mpga", "m4a", "wav", "webm"]:
            ext = "webm"

        client = Groq(api_key=settings.GROQ_API_KEY)
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as temp_file:
            temp_file.write(file_bytes)
            temp_path = temp_file.name

        try:
            with open(temp_path, "rb") as file:
                transcription = client.audio.transcriptions.create(
                    file=(f"audio.{ext}", file.read()),
                    model="whisper-large-v3",
                    response_format="json",
                )

            text = (
                transcription.get("text", "")
                if isinstance(transcription, dict)
                else getattr(transcription, "text", "")
            )
            if not text or not text.strip():
                raise ValueError("No speech detected.")
            return [
                {
                    "content": text.strip(),
                    "page_number": None,
                    "metadata": {"type": "audio"},
                }
            ]
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    def parse_youtube(self, url: str) -> list[dict]:
        """Fetch YouTube subtitles using yt_dlp."""
        ydl_opts = {
            "quiet": True,
            "skip_download": True,
            "writesubtitles": True,
            "writeautomaticsub": True,
            "subtitleslangs": ["en", "en-US", "en-GB"],
        }

        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                # Disable warning prints to keep terminal clean
                import logging

                ydl.logger = logging.getLogger("yt_dlp")
                ydl.logger.setLevel(logging.ERROR)

                info = ydl.extract_info(url, download=False)

                sub_url = None
                sub_format = None

                if "subtitles" in info and info["subtitles"]:
                    for lang in ["en", "en-US", "en-GB"]:
                        if lang in info["subtitles"]:
                            for sub in info["subtitles"][lang]:
                                if sub.get("ext") == "json3":
                                    sub_url = sub.get("url")
                                    sub_format = "json3"
                                    break
                            if not sub_url:
                                sub_url = info["subtitles"][lang][0].get("url")
                                sub_format = info["subtitles"][lang][0].get("ext")
                            break

                if (
                    not sub_url
                    and "automatic_captions" in info
                    and info["automatic_captions"]
                ):
                    for lang in ["en", "en-US", "en-GB", "en-orig"]:
                        if lang in info["automatic_captions"]:
                            for sub in info["automatic_captions"][lang]:
                                if sub.get("ext") == "json3":
                                    sub_url = sub.get("url")
                                    sub_format = "json3"
                                    break
                            if not sub_url:
                                sub_url = info["automatic_captions"][lang][0].get("url")
                                sub_format = info["automatic_captions"][lang][0].get(
                                    "ext"
                                )
                            break

                if not sub_url:
                    raise ValueError(
                        "No English subtitles or auto-captions exist for this video."
                    )

                with httpx.Client(follow_redirects=True, timeout=15) as client:
                    res = client.get(sub_url)
                    res.raise_for_status()

                full_text = ""

                if sub_format == "json3":
                    data = res.json()
                    events = data.get("events", [])
                    for event in events:
                        segs = event.get("segs", [])
                        for seg in segs:
                            text = seg.get("utf8", "").strip()
                            if text and text != "\n":
                                full_text += text + " "
                else:
                    # Fallback parsing for VTT
                    lines = res.text.split("\n")
                    clean_lines = []
                    for line in lines:
                        if (
                            not line.strip()
                            or line.strip().isdigit()
                            or "-->" in line
                            or line.startswith("WEBVTT")
                        ):
                            continue
                        line = re.sub(r"<[^>]+>", "", line)
                        clean_lines.append(line.strip())
                    full_text = " ".join(clean_lines)

                full_text = re.sub(r"\s+", " ", full_text).strip()

                if not full_text:
                    raise ValueError(
                        "Transcript was downloaded but contained no readable text."
                    )

                segments = []
                chunk_size = 2000
                for i in range(0, len(full_text), chunk_size):
                    segments.append(
                        {
                            "content": full_text[i : i + chunk_size],
                            "page_number": None,
                            "metadata": {"type": "youtube"},
                        }
                    )

                return segments

        except Exception as e:
            raise ValueError(f"Failed to fetch YouTube transcript: {str(e)}")

    def _extract_video_id(self, url: str) -> str:
        patterns = [
            r"(?:v=)([0-9A-Za-z_-]{11})",
            r"(?:youtu\.be/)([0-9A-Za-z_-]{11})",
            r"(?:embed/)([0-9A-Za-z_-]{11})",
        ]
        for p in patterns:
            match = re.search(p, url)
            if match:
                return match.group(1)
        return url
